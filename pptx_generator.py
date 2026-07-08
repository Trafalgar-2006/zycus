"""
pptx_generator.py — Monthly executive PowerPoint generator.

Reads live pipeline outputs (scores, mc, mc_v2, cluster, dag_info) when called
from run_all(). Falls back to parsing weekly .md files for the --monthly CLI flag.

Slide structure:
  1. Cover — Programme Health Overview
  2. Cross-Project Trend Analysis & Naive-vs-Model Baseline
  3. Emerging Risks & Root-Cause Clusters   ← real cluster data
  4. Deadline Probability & Forecast        ← real MC numbers
  5. Key Findings & Recommendations         ← data-derived
  6. Appendix — Methodology & Data Quality
"""
from __future__ import annotations

import re
from datetime import datetime
from pathlib import Path
from typing import Any

import config

# Colour palette (Zycus-inspired dark blue + accent)
_DARK_BLUE  = (0x1A, 0x2F, 0x5E)
_MID_BLUE   = (0x2E, 0x5E, 0xA8)
_WHITE      = (0xFF, 0xFF, 0xFF)
_RED        = (0xC0, 0x39, 0x2B)
_AMBER      = (0xE6, 0x7E, 0x22)
_GREEN      = (0x27, 0xAE, 0x60)
_LIGHT_GREY = (0xF2, 0xF2, 0xF2)

_RAG_RGB  = {"Red": _RED, "Amber": _AMBER, "Green": _GREEN}
_RAG_EMOJI = {"Red": "🔴", "Amber": "🟡", "Green": "🟢"}


def _rgb(t):
    from pptx.dml.color import RGBColor
    return RGBColor(*t)


def _parse_weekly_reports(weekly_dir: str = config.WEEKLY_OUTPUT_DIR) -> list[dict]:
    """Parse all weekly .md files — used as fallback when live_data is unavailable."""
    reports = []
    p = Path(weekly_dir)
    if not p.exists():
        return reports

    for md in sorted(p.glob("*.md")):
        text = md.read_text(encoding="utf-8")

        date_match = re.match(r"(\d{4}-\d{2}-\d{2})_(.*?)\.md", md.name)
        run_date   = date_match.group(1) if date_match else "Unknown"
        project    = date_match.group(2).replace("_", " ") if date_match else md.stem

        rag_match = re.search(r"Overall Status:.*?(Green|Amber|Red)", text)
        rag = rag_match.group(1) if rag_match else "Unknown"

        score_match = re.search(r"Project Score\s*\|\s*([\d.]+)", text)
        score = float(score_match.group(1)) if score_match else None

        # Prefer DAG v2 P(on-time) if present
        p_match_v2 = re.search(r"Dependency-aware DAG \(v2\).*?\*\*(\d+)%\*\*", text)
        p_match_v1 = re.search(r"Throughput.*?\*\*(\d+)%\*\*", text)
        p_on_time = (
            int(p_match_v2.group(1)) if p_match_v2 else
            int(p_match_v1.group(1)) if p_match_v1 else None
        )

        deadline_match = re.search(r"\*\*Deadline:\*\*\s*([\d-]+)", text)
        deadline = deadline_match.group(1) if deadline_match else "Unknown"

        cluster_match = re.search(r"\*\*(\d+ active tasks[^\*]+)\*\*", text)
        cluster_headline = cluster_match.group(1) if cluster_match else ""

        reports.append({
            "run_date":         run_date,
            "project":          project,
            "rag":              rag,
            "score":            score,
            "p_on_time":        p_on_time,
            "deadline":         deadline,
            "cluster_headline": cluster_headline,
            "source_file":      str(md),
            # Not available from MD parsing
            "mc":    {"p_on_time": p_on_time / 100 if p_on_time else None,
                      "deadline": deadline},
            "mc_v2": None,
            "cluster": {"headline": cluster_headline, "clusters": [], "n_at_risk_active": 0},
            "scores": {"project_score": score or 0, "rag": rag,
                       "task_summary": {"by_rag": {}, "by_status": {}},
                       "critical_path": {"total": 0, "by_rag": {}}},
            "dag_info": None,
        })

    return reports


def _add_slide(prs, layout_idx=6):
    layout = prs.slide_layouts[layout_idx]
    return prs.slides.add_slide(layout)


def _set_bg(slide, rgb_tuple):
    from pptx.dml.color import RGBColor
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(*rgb_tuple)


def _add_textbox(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=_WHITE, align="left"):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = RGBColor(*color)
    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    elif align == "right":
        p.alignment = PP_ALIGN.RIGHT
    return txBox


def _add_rag_box(slide, rag, left, top, size=0.9):
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    color = _RAG_RGB.get(rag, (0x99, 0x99, 0x99))
    shape = slide.shapes.add_shape(
        1,
        Inches(left), Inches(top), Inches(size), Inches(size)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*color)
    shape.line.fill.background()

    tf = shape.text_frame
    tf.paragraphs[0].text = rag[0]
    tf.paragraphs[0].font.size = Pt(28)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(*_WHITE)
    tf.paragraphs[0].alignment = __import__('pptx').enum.text.PP_ALIGN.CENTER


def generate_monthly_pptx(
    weekly_dir: str = config.WEEKLY_OUTPUT_DIR,
    output_dir: str = config.MONTHLY_OUTPUT_DIR,
    live_data: dict[str, dict[str, Any]] | None = None,
) -> Path:
    """
    Generate the monthly executive PowerPoint.

    Parameters
    ----------
    live_data : dict keyed by project name, each value contains:
        scores, mc, mc_v2, cluster, dag_info, df
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    month_str = datetime.now().strftime("%B %Y")
    date_str  = datetime.now().strftime("%Y-%m")

    # ── Resolve project data ─────────────────────────────────────────────────
    if live_data:
        # Use pipeline outputs directly — most accurate
        project_data: dict[str, dict] = live_data
    else:
        # Fallback: parse existing MD files
        parsed = _parse_weekly_reports(weekly_dir)
        # Keep only the latest run per project
        projects_tmp: dict[str, list] = {}
        for r in parsed:
            projects_tmp.setdefault(r["project"], []).append(r)
        project_data = {
            name: sorted(runs, key=lambda x: x["run_date"], reverse=True)[0]
            for name, runs in projects_tmp.items()
        }

    prs = Presentation()
    prs.slide_width  = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # ── SLIDE 1 — Cover ───────────────────────────────────────────────────────
    s1 = _add_slide(prs)
    _set_bg(s1, _DARK_BLUE)

    _add_textbox(s1, "ZYCUS", 0.5, 0.3, 12, 0.6, font_size=14, color=_MID_BLUE)
    _add_textbox(s1, "Programme Health Report", 0.5, 1.0, 12, 1.2,
                 font_size=44, bold=True, color=_WHITE, align="center")
    _add_textbox(s1, month_str, 0.5, 2.4, 12, 0.8,
                 font_size=26, color=(0xBB, 0xCC, 0xDD), align="center")
    _add_textbox(s1, "Prepared for Executive Leadership", 0.5, 3.1, 12, 0.5,
                 font_size=14, color=(0x99, 0xAA, 0xBB), align="center")

    x = 3.5
    for name, pdata in project_data.items():
        rag = pdata["scores"]["rag"]
        _add_rag_box(s1, rag, x, 4.2)
        _add_textbox(s1, name[:20], x - 0.1, 5.2, 1.8, 0.4,
                     font_size=10, color=_WHITE)
        x += 3.2

    _add_textbox(s1, f"Confidential | {datetime.now().strftime('%d %b %Y')}",
                 0.5, 7.0, 12, 0.4, font_size=10,
                 color=(0x77, 0x88, 0x99), align="center")

    # ── SLIDE 2 — Cross-Project Trend & Naive Baseline ───────────────────────
    s2 = _add_slide(prs)
    _set_bg(s2, (0xF8, 0xF9, 0xFA))

    _add_textbox(s2, "Cross-Project Health & Naive Baseline", 0.4, 0.2, 12, 0.7,
                 font_size=28, bold=True, color=_DARK_BLUE)
    _add_textbox(s2, f"Snapshot: {month_str}", 0.4, 0.85, 6, 0.4,
                 font_size=12, color=(0x55, 0x55, 0x55))

    y = 1.5
    scores_list = []
    for name, pdata in project_data.items():
        sc    = pdata["scores"]
        mc    = pdata["mc"]
        mc_v2 = pdata.get("mc_v2")
        rag   = sc["rag"]
        score = sc.get("project_score", 0)
        scores_list.append(score)

        # Project name + RAG badge
        _add_textbox(s2, f"{_RAG_EMOJI.get(rag, '⚪')} {name}",
                     0.4, y, 4.0, 0.45, font_size=14, bold=True, color=_DARK_BLUE)

        # Score bar
        bar_w = max(float(score) * 6, 0.05)
        bar = s2.shapes.add_shape(
            1, Inches(4.5), Inches(y + 0.05), Inches(bar_w), Inches(0.35)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = RGBColor(*_RAG_RGB.get(rag, (0x99, 0x99, 0x99)))
        bar.line.fill.background()

        # Score + P(on-time)
        p_v2  = (mc_v2 or {}).get("p_on_time")
        p_v1  = mc.get("p_on_time")
        p_str = (f"P(on-time) DAG: {round(p_v2*100)}%  |  v1: {round(p_v1*100)}%"
                 if p_v2 is not None and p_v1 is not None else
                 f"P(on-time): {round(p_v1*100)}%" if p_v1 is not None else "P(on-time): N/A")

        _add_textbox(s2, f"Score: {score:.2f}/1.00   {p_str}",
                     4.5, y + 0.42, 8.3, 0.35, font_size=11, color=(0x44, 0x44, 0x44))

        y += 1.35

    # Naive-vs-model footnote
    avg_score = sum(scores_list) / len(scores_list) if scores_list else 0
    red_count = sum(1 for p in project_data.values() if p["scores"]["rag"] == "Red")
    _add_textbox(
        s2,
        f"Portfolio avg score: {avg_score:.2f}/1.00 — {red_count}/{len(project_data)} project(s) Red.  "
        "Monte Carlo P(on-time) accounts for task velocity and variance "
        "— not just % complete. See Slide 4 for model comparison.",
        0.4, 6.3, 12.5, 0.9, font_size=12, color=_DARK_BLUE,
    )

    # ── SLIDE 3 — Emerging Risks & Root-Cause Clusters ───────────────────────
    s3 = _add_slide(prs)
    _set_bg(s3, _WHITE)

    _add_textbox(s3, "Emerging Risks & Root-Cause Analysis", 0.4, 0.2, 12, 0.7,
                 font_size=28, bold=True, color=_DARK_BLUE)

    # Risk bullets — built from real cluster + scores data
    bullets: list[str] = []
    for name, pdata in project_data.items():
        sc       = pdata["scores"]
        cluster  = pdata["cluster"]
        mc       = pdata["mc"]
        rag      = sc["rag"]
        n_red    = sc["task_summary"]["by_rag"].get("Red", 0)
        n_active = cluster.get("n_at_risk_active", 0)
        clusters = cluster.get("clusters", [])
        top_owner = clusters[0]["top_owner"] if clusters else "various owners"
        top_n     = clusters[0]["task_count"] if clusters else 0
        p_on      = mc.get("p_on_time")
        p_str     = f"{round(p_on*100)}% P(on-time)" if p_on is not None else "forecast N/A"
        emoji     = _RAG_EMOJI.get(rag, "⚪")

        bullets.append(
            f"{emoji}  {name} — {rag.upper()}: {n_active} active at-risk tasks, "
            f"{n_red} Red. Top cluster: {top_n} tasks under {top_owner}. {p_str}."
        )

    # Cross-project systemic insight
    red_projects = [n for n, p in project_data.items() if p["scores"]["rag"] == "Red"]
    if red_projects:
        bullets.append(
            f"⚠️  {len(red_projects)}/{len(project_data)} project(s) Red. "
            "Risk clusters show owner concentration — potential shared resource bottleneck."
        )
    else:
        bullets.append(
            "ℹ️  No projects currently Red. Continue monitoring Amber items and critical-path tasks."
        )

    y_b = 1.1
    for bullet in bullets[:5]:   # cap at 5 to avoid overflow
        _add_textbox(s3, bullet, 0.5, y_b, 12.3, 0.75, font_size=13, color=(0x22, 0x22, 0x22))
        y_b += 0.88

    # Risk concentration table header
    _add_textbox(s3, "Risk Concentration by Project", 0.4, 5.05, 12, 0.45,
                 font_size=14, bold=True, color=_DARK_BLUE)

    headers   = ["Project", "At-Risk Active", "Top Owner / Cluster", "P(on-time)"]
    col_w     = [3.0, 2.2, 5.3, 2.0]
    x_starts  = [0.4, 3.5, 5.8, 11.2]

    for hdr, w, x in zip(headers, col_w, x_starts):
        hb = s3.shapes.add_shape(1, Inches(x), Inches(5.55), Inches(w), Inches(0.4))
        hb.fill.solid()
        hb.fill.fore_color.rgb = RGBColor(*_DARK_BLUE)
        hb.line.fill.background()
        _add_textbox(s3, hdr, x + 0.05, 5.58, w - 0.1, 0.35,
                     font_size=11, bold=True, color=_WHITE)

    y_row = 6.02
    for name, pdata in project_data.items():
        mc       = pdata["mc"]
        mc_v2    = pdata.get("mc_v2")
        cluster  = pdata["cluster"]
        clusters = cluster.get("clusters", [])
        top_owner = clusters[0]["top_owner"] if clusters else "—"
        top_n     = clusters[0]["task_count"] if clusters else 0
        p_v2     = (mc_v2 or {}).get("p_on_time")
        p_v1     = mc.get("p_on_time")
        p_str    = (f"{round(p_v2*100)}% (DAG)"
                    if p_v2 is not None else
                    f"{round(p_v1*100)}%" if p_v1 is not None else "N/A")

        row_vals = [
            name[:28],
            str(cluster.get("n_at_risk_active", 0)),
            f"{top_n} tasks / {top_owner}"[:40],
            p_str,
        ]
        for val, w, x in zip(row_vals, col_w, x_starts):
            _add_textbox(s3, val, x + 0.05, y_row, w - 0.1, 0.38,
                         font_size=11, color=(0x22, 0x22, 0x22))
        y_row += 0.42

    # ── SLIDE 4 — Deadline Probability Forecast ───────────────────────────────
    s4 = _add_slide(prs)
    _set_bg(s4, (0xF0, 0xF4, 0xF8))

    _add_textbox(s4, "Deadline Probability Forecast", 0.4, 0.2, 12, 0.7,
                 font_size=28, bold=True, color=_DARK_BLUE)
    _add_textbox(
        s4,
        "Monte Carlo: 10,000 scenarios per project using historical task duration ratios "
        "(actual/planned) from completed tasks. "
        "v2 (DAG) propagates durations through predecessor graph — respects sequencing structure.",
        0.4, 0.95, 12.5, 0.65, font_size=12, color=(0x44, 0x44, 0x55),
    )

    y = 1.85
    for name, pdata in project_data.items():
        mc    = pdata["mc"]
        mc_v2 = pdata.get("mc_v2")
        rag   = pdata["scores"]["rag"]

        deadline = mc.get("deadline", "Unknown")
        p_on_v1  = mc.get("p_on_time")
        p_on_v2  = (mc_v2 or {}).get("p_on_time")
        p_1w     = mc.get("p_slip_1w")
        p_2w     = mc.get("p_slip_2w_plus")

        # Always use numeric comparisons — never parse strings
        p_v1_str = f"{round(p_on_v1*100)}%" if p_on_v1 is not None else "N/A"
        p_v2_str = f"{round(p_on_v2*100)}%" if p_on_v2 is not None else "—"
        p_1w_str = f"{round(p_1w*100)}%"    if p_1w  is not None else "N/A"
        p_2w_str = f"{round(p_2w*100)}%"    if p_2w  is not None else "N/A"

        p_color   = _GREEN if (p_on_v1 or 0) > 0.60 else (_AMBER if (p_on_v1 or 0) > 0.25 else _RED)
        p2w_color = _RED   if (p_2w    or 0) > 0.25 else _AMBER

        box = s4.shapes.add_shape(1, Inches(0.4), Inches(y), Inches(12.5), Inches(1.85))
        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(*_WHITE)
        box.line.color.rgb      = RGBColor(*_RAG_RGB.get(rag, (0x99, 0x99, 0x99)))

        _add_textbox(s4, name, 0.6, y + 0.10, 4.0, 0.45, font_size=16, bold=True, color=_DARK_BLUE)
        _add_textbox(s4, f"Deadline: {deadline}", 0.6, y + 0.60, 4.0, 0.35,
                     font_size=12, color=(0x44, 0x44, 0x44))

        # v1 / v2 comparison label
        if p_on_v2 is not None:
            diff_pp = round(p_on_v2 * 100) - round(p_on_v1 * 100)
            label   = (f"v1 (throughput): {p_v1_str}  →  v2 (DAG): {p_v2_str}  "
                       f"({diff_pp:+d} pp from dependency structure)")
        else:
            label = f"v1 (throughput): {p_v1_str}  |  v2 skipped (low predecessor coverage)"
        _add_textbox(s4, label, 0.6, y + 1.05, 11.7, 0.45, font_size=11, color=(0x44, 0x44, 0x44))

        # Three metric boxes
        metrics = [
            ("P(on time)",    p_v1_str, p_color),
            ("P(slip ≤1wk)",  p_1w_str, _AMBER),
            ("P(slip 2wk+)",  p_2w_str, p2w_color),
        ]
        mx = 5.2
        for label_m, val, col in metrics:
            _add_textbox(s4, label_m, mx, y + 0.10, 2.2, 0.35, font_size=11,
                         color=(0x55, 0x55, 0x55))
            _add_textbox(s4, val, mx, y + 0.45, 2.2, 0.65, font_size=28,
                         bold=True, color=col)
            mx += 2.45

        y += 2.1

    # ── SLIDE 5 — Key Findings & Recommendations ─────────────────────────────
    s5 = _add_slide(prs)
    _set_bg(s5, _DARK_BLUE)

    _add_textbox(s5, "Key Findings & Recommendations", 0.4, 0.2, 12, 0.7,
                 font_size=28, bold=True, color=_WHITE)

    # Findings derived from live data
    findings: list[tuple[str, str]] = []

    # Finding 1: MC headline — most at-risk project
    most_risk = max(project_data.items(), key=lambda x: x[1]["scores"]["project_score"])
    mr_name, mr_data = most_risk
    mr_mc   = mr_data["mc"]
    mr_mc2  = mr_data.get("mc_v2")
    mr_p    = mr_mc.get("p_on_time")
    mr_p2   = (mr_mc2 or {}).get("p_on_time")
    p_detail = (
        f"v1: {round(mr_p*100)}%, v2 (DAG): {round(mr_p2*100)}%"
        if mr_p2 is not None and mr_p is not None else
        f"{round(mr_p*100)}% (throughput model)" if mr_p is not None else "N/A"
    )
    findings.append((
        "🔍 Deadline Risk",
        f"{mr_name} is the highest-risk project (score {mr_data['scores']['project_score']:.3f}). "
        f"P(on-time) {p_detail}. Deadline: {mr_mc.get('deadline', 'Unknown')}.",
    ))

    # Finding 2: Critical path analysis per project
    for name, pdata in project_data.items():
        sc   = pdata["scores"]
        cp   = sc.get("critical_path", {})
        dag  = pdata.get("dag_info")
        red_cp = cp.get("by_rag", {}).get("Red", 0)
        cp_total = cp.get("total", 0)
        dag_cp_n = len(dag.get("critical_path", [])) if dag else None
        dag_cov  = (dag or {}).get("coverage", {}).get("pct_coverage", 0)
        if cp_total > 0:
            cp_detail = f"{red_cp}/{cp_total} PM-flagged CP tasks are Red"
            if dag_cp_n is not None and dag_cov >= 0.50:
                cp_detail += f"; graph-computed CP = {dag_cp_n} tasks ({dag_cov:.0%} predecessor coverage)"
            findings.append(("🔍 Critical Path", f"{name}: {cp_detail}."))

    # Finding 3: Risk cluster concentration
    for name, pdata in project_data.items():
        cluster  = pdata["cluster"]
        clusters = cluster.get("clusters", [])
        if clusters:
            top = clusters[0]
            findings.append((
                "🔍 Risk Cluster",
                f"{name}: {cluster.get('headline', '')} "
                f"Top cluster — {top['task_count']} tasks, owner: {top['top_owner']}, "
                f"{top['red_tasks']} Red.",
            ))

    y_f = 1.1
    for title, body in findings[:3]:
        _add_textbox(s5, title, 0.5, y_f, 2.9, 0.35, font_size=12, bold=True,
                     color=(0xAA, 0xCC, 0xFF))
        _add_textbox(s5, body, 3.5, y_f, 9.3, 0.55, font_size=12, color=_WHITE)
        y_f += 0.78

    # Actions derived from cluster analysis
    y_f += 0.15
    _add_textbox(s5, "Recommended Actions", 0.4, y_f, 12, 0.45,
                 font_size=16, bold=True, color=(0xAA, 0xCC, 0xFF))
    y_f += 0.55

    recs: list[tuple[str, str]] = []
    for name, pdata in project_data.items():
        cluster  = pdata["cluster"]
        clusters = cluster.get("clusters", [])
        n_at_risk = cluster.get("n_at_risk_active", 0)
        top_owner = clusters[0]["top_owner"] if clusters else "assigned owners"
        recs.append((
            f"✅ {name}",
            f"Triage {n_at_risk} active at-risk tasks with {top_owner}; "
            "escalate any CP tasks that are Red or slipping.",
        ))
    recs.append((
        "✅ Data Quality",
        "Ensure Predecessors column is populated in MS Project exports "
        "to enable dependency-aware simulation on all projects.",
    ))

    for title, body in recs[:3]:
        _add_textbox(s5, title, 0.5, y_f, 2.9, 0.35, font_size=12, bold=True, color=_GREEN)
        _add_textbox(s5, body, 3.5, y_f, 9.3, 0.55, font_size=12, color=_WHITE)
        y_f += 0.72

    # ── SLIDE 6 — Appendix: Methodology & Data Quality ───────────────────────
    s6 = _add_slide(prs)
    _set_bg(s6, _WHITE)

    _add_textbox(s6, "Appendix — Methodology & Data Quality", 0.4, 0.2, 12, 0.65,
                 font_size=24, bold=True, color=_DARK_BLUE)

    # Feature importance from live data (first project's shap_info, if available)
    fi_note = ""
    for pdata in project_data.values():
        shap = pdata.get("shap_info", {})
        top  = shap.get("top_features", [])
        if top:
            fi_note = (
                "ML Model top features (GBT, verified this run): "
                + ", ".join(f"{f['feature']} ({f['mean_abs_shap']})" for f in top[:3])
                + ". "
            )
            break

    method_text = (
        "RAG Score = 0.60 × Forward Risk + 0.40 × Historical Slip. "
        "Forward Risk: fraction of critical-path active tasks that are Red (×1.0) or Amber (×0.5). "
        "Falls back to all active tasks if critical count <5.\n"
        "Historical Slip: median(|variance_days|) for late completed tasks, normalised against 30-day ceiling. "
        "Variance column sign used (verified: 196/199 rows correct). Magnitude treated as unreliable for parent rows.\n"
        + fi_note
        + "RAG is manually set by PMs — model learns PM judgment patterns, not a formula. "
        "Fit/interpret only — 2 projects, no holdout split.\n"
        "Monte Carlo v1: throughput rate (tasks/week) × log-normal duration-ratio distribution. "
        "MC v2 (DAG): discrete-event simulation propagating durations topologically through predecessor graph. "
        "v2 runs only when predecessor coverage ≥50%; otherwise v1 is sole forecast.\n"
        "Clustering: TF-IDF + KMeans on available fields (Area, Phase, Owner, Task Name). "
        "Sparsity-aware: fields missing in >98% of rows excluded; caveat lists actual population %.\n"
        "Data confidence tiers — High: weighted null rate <20%; Medium: 20–45%; Low: >45%. "
        "Predecessor coverage weighted 0.4; planned_days + pct_complete 0.3 each."
    )
    _add_textbox(s6, method_text, 0.4, 1.0, 12.5, 5.8,
                 font_size=11, color=(0x22, 0x22, 0x33))

    # ── Save ──────────────────────────────────────────────────────────────────
    out_dir = Path(output_dir)
    out_dir.mkdir(parents=True, exist_ok=True)
    out_path = out_dir / f"Executive_Report_{date_str}.pptx"
    prs.save(str(out_path))
    return out_path
